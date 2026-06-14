# Builds sokol-ads.pptx: bg renders + !!-named morph objects (logo, LED screens, bar, wing).
import io
from PIL import Image
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

def bg_jpeg(path):
    """Compress a fresh bg PNG to in-memory JPEG q88 (no stale .jpg on disk)."""
    im = Image.open(path).convert("RGB")
    buf = io.BytesIO()
    im.save(buf, "JPEG", quality=88)
    buf.seek(0)
    return buf

PX = 6350  # EMU per px at 1920px -> 13.333in
ROOT = r"C:\Users\user\presentation"
S = ROOT + r"\ads\shots"
O = ROOT + r"\ads\obj"

def E(px): return Emu(int(round(px * PX)))

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]

NSMAP = ('xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
         'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
         'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
         'xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main"')
MORPH_XML = (f'<mc:AlternateContent {NSMAP}>'
             '<mc:Choice Requires="p159">'
             '<p:transition spd="slow" p14:dur="1000"><p159:morph option="byObject"/></p:transition>'
             '</mc:Choice>'
             '<mc:Fallback><p:transition spd="slow"><p:fade/></p:transition></mc:Fallback>'
             '</mc:AlternateContent>')
FADE_XML = (f'<p:transition {NSMAP} spd="med" p14:dur="500"><p:fade/></p:transition>')

def set_transition(slide, xml):
    el = etree.fromstring(xml)
    sld = slide._element
    MC = '{http://schemas.openxmlformats.org/markup-compatibility/2006}AlternateContent'
    for tag in (qn('p:transition'), MC):
        for old in sld.findall(tag):
            sld.remove(old)
    sld.find(qn('p:clrMapOvr')).addnext(el)

# ---- geometry, deck px space (mirrors GEOM in sokol-ads.html) ----
# logo.png rendered from obj/logo.svg (vector), aspect 121:35 = 3.457
LOGO = {1: (710, 176, 500, 145), 2: (60, 38, 150, 43), 3: (60, 38, 150, 43),
        4: (60, 38, 150, 43), 5: (760, 311, 400, 116)}
WING = {1: (1210, 500, 760, 560), 2: (-340, -220, 760, 560), 3: (-340, -240, 760, 560),
        4: (-340, -240, 760, 560), 5: (1330, -200, 760, 560)}
# !!led = text-only lettering on the blacked-out screens: s3 totem (portrait) -> s4 MEGA panel
LED = {3: ("obj-ledtxt3.png", (325, 313, 350, 525)),
       4: ("obj-ledtxt4.png", (248, 148, 900, 630))}

def put(slide, img, box, name):
    x, y, w, h = box
    pic = slide.shapes.add_picture(img, E(x), E(y), E(w), E(h))
    pic.name = name
    return pic

for n in range(1, 6):
    sl = prs.slides.add_slide(BLANK)
    sl.shapes.add_picture(bg_jpeg(rf"{S}\bg{n}.png"), 0, 0, prs.slide_width, prs.slide_height)
    put(sl, rf"{S}\obj-wing.png", WING[n], "!!wing")
    if n in LED:
        img, box = LED[n]
        put(sl, rf"{S}\{img}", box, "!!led")
    put(sl, rf"{O}\logo.png", LOGO[n], "!!logo")
    set_transition(sl, MORPH_XML if n > 1 else FADE_XML)

out = rf"{ROOT}\sokol-ads.pptx"
prs.save(out)
import os
print("slides:", len(prs.slides._sldIdLst), "->", out, os.path.getsize(out) // 1024, "KB")
